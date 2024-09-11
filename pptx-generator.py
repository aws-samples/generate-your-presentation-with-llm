import sys
import os
import json
import uuid
import time
from src.template_mapping import template_aws1
import streamlit as st
from datetime import datetime
import subprocess
from PIL import Image
import glob
import shlex
import ast
from pptx import Presentation, util

from src.utils import check_password
from src.utils import generate_text
from src.utils import generate_bedrock_image
from src.utils import check_text_generation_consistency
from src.utils import invoke_llm_text
from src.utils import is_valid_text_gen_json
from src.utils import validate_slide_json

from src.prompt import create_initial_prompt
from src.prompt import moderation_prompt
from src.prompt import agenda_prompt
from src.prompt import example_json

# USERNAME AND PASSWORD
if not st.session_state.get("glb_username"):
    st.session_state["glb_username"] = st.secrets["my_secrets"]["glb_username"]
if not st.session_state.get("glb_pwd"):
    st.session_state["glb_pwd"] = st.secrets["my_secrets"]["glb_pwd"]

# FIRST STREAMLIT COMMAND
st.set_page_config(layout="wide")

# Design move app further up and remove top padding
st.markdown('''<style>.css-1egvi7u {margin-top: -3rem;}</style>''',
            unsafe_allow_html=True)
# Design change st.Audio to fixed height of 45 pixels
st.markdown('''<style>.stAudio {height: 45px;}</style>''',
            unsafe_allow_html=True)
# Design change hyperlink href link color
st.markdown('''<style>.css-v37k9u a {color: #ff4c4b;}</style>''',
            unsafe_allow_html=True)  # darkmode
st.markdown('''<style>.css-nlntq9 a {color: #ff4c4b;}</style>''',
            unsafe_allow_html=True)  # lightmode

st.markdown(
                """
                <style>

                div.row-widget.stRadio > div > label {
                    background-color: rgb(19, 23, 32);
                    padding: 5px 10px;
                    margin-right: 5px;
                    margin-bottom: 5px;
                    border-radius: 10px;
                    border: 1px solid rgba(250, 250, 250, 0.2);
                    display: inline-flex;
                    align-items: center;
                    justify-content: left;
                    min-width: 120px;
                    text-align: left;
                }
                
                div.row-widget.stRadio > div > label:hover {
                    border-color: #ff4c4b;
                }
                
                div.row-widget.stRadio > div > label > div:first-child > div {
                    background-color: transparent !important;
                    border-color: transparent !important;
                }
                
                div.row-widget.stRadio > div > label.stRadio > div:first-child > div:after {
                    content: '';
                }
                
                </style>
                """, unsafe_allow_html=True
            )

if not check_password("PPTX generator App"):
    # need to login first
    st.stop()

# MAIN PAGE        
st.title('Generate your presentation with Amazon Bedrock!')

if not st.session_state.get("TOPIC"):
    st.session_state["TOPIC"] = 'Benefits of cloud computing with Amazon Web Services'
if not st.session_state.get("TOPIC_FROM_TEXT"):
    st.session_state["TOPIC_FROM_TEXT"] = 'Benefits of cloud computing with Amazon Web Services'
if not st.session_state.get("BKG_PROMPT"):
    st.session_state["BKG_PROMPT"] = "digital presentation wallpaper, dark blue tone, uniform color, corner gradient towards orange"
if not st.session_state.get("your_full_name"):
    st.session_state["your_full_name"] = 'J. Doe'
if not st.session_state.get("your_contact_info"):
    st.session_state["your_contact_info"] = "j.doe@anycompany.com"
if not st.session_state.get("your_title"):
    st.session_state["your_title"] = "Chief Presentation Officer"
if not st.session_state.get("your_company"):
    st.session_state["your_company"] = "AnyCompany"

cwd = os.getcwd() 

# SET TO TRUE ONLY FOR TEMPLATE DEBUG PURPOSE
characterize_template = False
if characterize_template:
    prs = Presentation(cwd+"/templates/pptx_base_template.pptx")

    st.write("")

    for idx_slide, islide in enumerate(prs.slide_layouts):

        # Pick slide layout from master
        st.write("Slide "+str(idx_slide+1)+" -- idx: "+str(idx_slide))
        slide_layout = prs.slide_layouts[idx_slide]
        slide = prs.slides.add_slide(slide_layout)
        
        # Modify shapes
        for i,shape in enumerate(slide.placeholders):
            st.write(str('\tPlaceholder Index: %d - Name: %s' % (shape.placeholder_format.idx, shape.name)))
    
    raise ValueError("Finished enumerating the master slides placeholders")

inputs_col1, inputs_col2, inputs_col3 = st.columns(3)

with inputs_col1:
        
    st.session_state["selected_LLM"] = st.selectbox('Choose Language Model', ('Claude 3 Haiku', 'Claude 3 Sonnet'), index=1, key="LLM")
    if st.session_state["selected_LLM"] == 'Claude 3 Haiku':
        st.session_state["chosen_LLM"] = "anthropic.claude-3-haiku-20240307-v1:0"
        st.session_state["LLM_input_token_price"] =  0.00025/1e3 # us-east-1
        st.session_state["LLM_output_token_price"] = 0.00125/1e3 # us-east-1
    elif st.session_state["selected_LLM"] == 'Claude 3 Sonnet':
        st.session_state["chosen_LLM"] = "anthropic.claude-3-sonnet-20240229-v1:0"
        st.session_state["LLM_input_token_price"] =  0.00300/1e3 # us-east-1
        st.session_state["LLM_output_token_price"] = 0.01500/1e3 # us-east-1

    st.session_state["TOPIC"] = st.text_area("Insert your topic of choice", st.session_state["TOPIC_FROM_TEXT"], key="topic")
    st.session_state["N_SLIDES"] = st.slider("Preferred number of slides", min_value=5, max_value=15, value=6, step=1, format="%i", key="slides")

    st.session_state["create_agenda_checkbox"] = st.checkbox('Add agenda slide', value=True)
    st.session_state["create_thankyou_checkbox"] = st.checkbox('Add thank you slide', value=True)

    st.session_state["selected_generate_bkg"] = st.checkbox('Generate custom background image', value=True)
    
    if not st.session_state["selected_generate_bkg"]:
        st.session_state["generate_bkg"] = False
        st.session_state["gen_bkg_price_cents"] = 0
    else:
        st.session_state["generate_bkg"] = True
        st.session_state["generate_bkg_prompt"] = st.text_area("Background image prompt", st.session_state["BKG_PROMPT"], key="bkg_prompt", disabled=not st.session_state["selected_generate_bkg"])
        st.session_state["gen_bkg_price_cents"] = 1

    st.session_state["selected_generate_images"] = st.selectbox('Generate images within Slides', ('Do not generate', 'Low resolution', 'High resolution'), index=1, key="bkg_gen")
    if st.session_state["selected_generate_images"] == 'Do not generate':
        st.session_state["generate_images"] = False
        st.session_state["high_res_images"] = False
        st.session_state["gen_images_price_cents"] = 0
    else:
        st.session_state["generate_images"] = True
        if st.session_state["selected_generate_images"] == 'Low resolution':
            st.session_state["high_res_images"] = False
            st.session_state["gen_images_price_cents"] = 0.8
        else:
            st.session_state["high_res_images"] = True
            st.session_state["gen_images_price_cents"] = 1

    st.session_state["generate_thumbnails"] = st.checkbox('Generate slide thumbnails', value=True)

    if st.session_state["generate_thumbnails"]:
        st.warning('The thumbnails are not reflecting the exact layout of the slides (font, text alignment, images are missing)', icon="⚠️")

with inputs_col2:
    st.session_state["customize_contact_info"] = st.checkbox('Customize Contact Information', value=True)
    if st.session_state["customize_contact_info"]:
        st.session_state["your_full_name"] = st.text_input("Your Full Name", st.session_state["your_full_name"], key="your_full_name1", disabled=not st.session_state["customize_contact_info"])
        st.session_state["your_contact_info"] = st.text_input("Your Contact", st.session_state["your_contact_info"], key="your_contact_info1", disabled=not st.session_state["customize_contact_info"])
        st.session_state["your_title"] = st.text_input("Your Job Title", st.session_state["your_title"], key="your_title1", disabled=not st.session_state["customize_contact_info"])
        st.session_state["your_company"] = st.text_input("Your Company", st.session_state["your_company"], key="your_company1", disabled=not st.session_state["customize_contact_info"])

if st.button('Generate presentation', key="create_presentation"):
    st.session_state["n_input_tokens"] = 0
    st.session_state["n_output_tokens"] = 0
    st.session_state["n_gen_images"] = 0
    start_time = datetime.now()

    print("\nSTARTING PRESENTATION GENERATION!")
    st.session_state["content_allowed"] = False

    # Load Template Slides Formats
    st.session_state["slides_format_json"] = template_aws1(high_res_images = st.session_state["high_res_images"])
    
    moderate_request_response, usage = invoke_llm_text(moderation_prompt(TOPIC=st.session_state["TOPIC"]))
    st.session_state["n_input_tokens"] = st.session_state["n_input_tokens"] + usage['input_tokens']
    st.session_state["n_output_tokens"] = st.session_state["n_output_tokens"] + usage['output_tokens']
    
    try:
        st.session_state["content_allowed"] = ast.literal_eval( (moderate_request_response[0])["text"] )["content_allowed"]
    except ValueError:
        # need to implement proper retry
        print("content_allowed ERROR, skipping")
        
    print("content_allowed",st.session_state["content_allowed"])

    if st.session_state["content_allowed"].capitalize():

        initial_prompt = create_initial_prompt(N_SLIDES=st.session_state["N_SLIDES"], TOPIC=st.session_state["TOPIC"])

        generated_n_slides_is_consistent = False
        max_generation_attempts = 2
        generation_attempts = 0

        st.warning('Generating presentation... The application will try to self-heal in case of errors, click the button again if it fails to do so 😉', icon="🚨")

        while (not generated_n_slides_is_consistent):

            text_gen_result, usage = generate_text(prompt=initial_prompt, N_SLIDES=st.session_state["N_SLIDES"], model_id = st.session_state["chosen_LLM"])
            st.session_state["n_input_tokens"] = st.session_state["n_input_tokens"] + usage['input_tokens']
            st.session_state["n_output_tokens"] = st.session_state["n_output_tokens"] + usage['output_tokens']

            if len(text_gen_result) != 1:
                st.write("There was a problem with the answer from Claude, try again")
                sys.exit()

            print("RAW RESULT:")
            print(text_gen_result)

            raw_generated_json = text_gen_result[0]["text"]
            print("MANIPULATION ATTEMPT RESULT:")
            print("RAW GENERATED JSON:",raw_generated_json)

            is_valid_json_content = is_valid_text_gen_json(raw_json=raw_generated_json)
            print("is_valid_json_content",is_valid_json_content)

            gen_result_fix_attempts = 0
            max_gen_result_fix_attempts = 2
            while not is_valid_json_content and gen_result_fix_attempts < max_gen_result_fix_attempts:
                gen_result_fix_attempts = gen_result_fix_attempts+1
                print("FIXING BROKEN INITIALLY GENERATED JSON!")
                raw_gen_result_fixed, usage = generate_text(prompt=("Fix the broken JSON below:\n- Use empty strings instead of Null values\n-Escape quotes\n- Return a valid JSON format\n\nVery important: Skip the preamble\n\Here is the broken JSON:\n"+str(raw_generated_json)), N_SLIDES=st.session_state["N_SLIDES"], model_id = st.session_state["chosen_LLM"])
                st.session_state["n_input_tokens"] = st.session_state["n_input_tokens"] + usage['input_tokens']
                st.session_state["n_output_tokens"] = st.session_state["n_output_tokens"] + usage['output_tokens']
                tmp_raw_generated_json = raw_gen_result_fixed[0]["text"]
                print("tmp_raw_generated_json",tmp_raw_generated_json)
                is_valid_json_content = is_valid_text_gen_json(raw_json=tmp_raw_generated_json)
                print("attempt",gen_result_fix_attempts,"is_valid_json_content",is_valid_json_content)
                if is_valid_json_content:
                    raw_generated_json = tmp_raw_generated_json
                else:
                    st.warning('Generated slides JSON contains bugs, fixing it...', icon="🚨")
            
            st.session_state["valid_generation"] = True
            try:
                validated_json_content = ast.literal_eval(raw_generated_json)
            except ValueError:
                try:
                    # print("TRYING WITH: ",(raw_generated_json+"\"}]}"))
                    # bad hack to try to force valid json in case of truncation
                    validated_json_content = ast.literal_eval(raw_generated_json+"\"}]}")
                except ValueError:
                    st.error('Errors encountered in the generation, please try again', icon="🚨")
                    st.session_state["valid_generation"] = False
            
            print(type(validated_json_content))
            json_content_list = validated_json_content["slides"]
            print("Generated slides content:")
            print(json_content_list)

            generated_n_slides, generated_n_slides_is_consistent = check_text_generation_consistency(slides_list=json_content_list, N_SLIDES=st.session_state["N_SLIDES"])
            generation_attempts = generation_attempts+1
            if not generated_n_slides_is_consistent:
                st.warning("INCONSISTENT NUMBER OF SLIDES! "+str(st.session_state["N_SLIDES"])+" requested, "+str(generated_n_slides)+" generated", icon="🚨")
                if generation_attempts < max_generation_attempts:
                    st.write("WAITING 3 seconds and trying again")
                    time.sleep(3) # nosemgrep: arbitrary-sleep, waiting 3 seconds to avoid throttling
                else:
                    st.write("MOVING FORWARD ANYWAY after "+str(max_generation_attempts)+" attempts")
                    break

        if st.session_state["valid_generation"]:
            prs = Presentation(cwd+"/templates/pptx_base_template.pptx")

            st.write("")

            if st.session_state["selected_generate_bkg"]:
                generate_bedrock_image(img_prompt=st.session_state["generate_bkg_prompt"], current_slide_format_json={"image_height": 768, "image_width": 1152}, 
                # image_placeholder=image_placeholder, 
                cwd=cwd, bkg="_bkg")
                st.session_state["n_gen_images"] = st.session_state["n_gen_images"]+1

            # TRY TO FIX INDIVIDUAL SLIDE JSONS
            raw_json_content_list = json_content_list
            for i_json_slide, json_slide in enumerate(raw_json_content_list):
                
                is_valid_json_slide = validate_slide_json(slide_json=json_slide)
                print("is_valid_json_slide",is_valid_json_slide)

                slide_json_fix_attempts = 0
                max_slide_json_fix_attempts = 3
                while not is_valid_json_slide and slide_json_fix_attempts < max_slide_json_fix_attempts:
                    slide_json_fix_attempts = slide_json_fix_attempts+1
                    print("FIXING BROKEN SLIDE JSON!")
                    tmp_result_slide_fixed, usage = generate_text(prompt=(example_json()+"""\nFix and add the missing fields to the broken JSON below, according to the slide examples above. Return a valid JSON format."""+str(json_slide)), N_SLIDES=st.session_state["N_SLIDES"], model_id = st.session_state["chosen_LLM"])
                    st.session_state["n_input_tokens"] = st.session_state["n_input_tokens"] + usage['input_tokens']
                    st.session_state["n_output_tokens"] = st.session_state["n_output_tokens"] + usage['output_tokens']
                    raw_generated_slide_json = tmp_result_slide_fixed[0]["text"]
                    print("raw_generated_slide_json",raw_generated_slide_json)
                    is_valid_json_slide = validate_slide_json(slide_json=raw_generated_slide_json)
                    print("attempt",slide_json_fix_attempts,"is_valid_json_slide",is_valid_json_slide)
                    if is_valid_json_slide:
                        json_slide = ast.literal_eval(raw_generated_slide_json)

                json_content_list[i_json_slide] = json_slide


            # NOW CREATE SLIDES FROM VALIDATED JSONS
            tot_n_slides = generated_n_slides + 1*st.session_state["create_agenda_checkbox"] + 1*st.session_state["create_thankyou_checkbox"]
            outputs_col1 = [None] * tot_n_slides
            outputs_col2 = [None] * tot_n_slides
            outputs_col3 = [None] * tot_n_slides

            i_slide_col = 0
            raw_json_content_list = json_content_list
            for validated_slide_json_content in raw_json_content_list:
                
                outputs_col1[i_slide_col], outputs_col2[i_slide_col], outputs_col3[i_slide_col] = st.columns(3)

                with outputs_col1[i_slide_col]:

                    # Pick slide layout from master
                    current_slide_format = validated_slide_json_content["slideFormat"]
                    st.write("Slide "+str(i_slide_col+1)+" - "+current_slide_format)
                    st.write(validated_slide_json_content)
                    try:
                        current_slide_format_json = st.session_state["slides_format_json"][current_slide_format]
                    except ValueError:
                        #force fallback to "Slide with image and text"
                        current_slide_format_json = "Slide with image and text"
                    layout_slide_idx = current_slide_format_json["layout_slide"]
                    slide_layout = prs.slide_layouts[layout_slide_idx]
                    slide = prs.slides.add_slide(slide_layout)

                    notes_slide = slide.notes_slide
                    text_frame = notes_slide.notes_text_frame
                    text_frame.text = json.dumps(validated_slide_json_content, indent=4)

                    # Check shapes
                    print("layout_slide_idx",layout_slide_idx)
                    print("current_slide_format",current_slide_format)
                    for i,shape in enumerate(slide.placeholders):
                        print(str('\tPlaceholder Index: %d - Name: %s' % (shape.placeholder_format.idx, shape.name)))

                    if current_slide_format == "Title page":
                        # Slide title
                        title = slide.placeholders[current_slide_format_json["title_placeholder"]]
                        title.text = validated_slide_json_content["title"]
                        # Slide subtitle
                        subtitle = slide.placeholders[current_slide_format_json["subtitle_placeholder"]]
                        subtitle.text = validated_slide_json_content["subtitle"]
                        # Slide Author Name
                        if st.session_state["customize_contact_info"]:
                            full_name = slide.placeholders[current_slide_format_json["full_name_placeholder"]] 
                            full_name.text = st.session_state["your_full_name"]
                        # Slide Author Title
                        if st.session_state["customize_contact_info"]:
                            job_title = slide.placeholders[current_slide_format_json["job_title_placeholder"]] 
                            job_title.text = st.session_state["your_title"]+"\n"+st.session_state["your_company"]
                        # Customize background picture
                        left = top = util.Inches(0)
                        if st.session_state["selected_generate_bkg"]:
                            bkg_img_path = "/tmp/test_image_bkg.jpg"
                        else:
                            bkg_img_path = "/templates/default_bkg.jpg"
                        pic = slide.shapes.add_picture(cwd+bkg_img_path, left, top, width=prs.slide_width, height=prs.slide_height)
                        cursor_sp = slide.shapes[0]._element
                        cursor_sp.addprevious(pic._element)

                        # Add original input before  generated output
                        text_frame.text = "Input topic/text:\n" + st.session_state.get("TOPIC") + "\n\nGenerated content:\n" + text_frame.text

                    elif current_slide_format == "Slide with bullet points":
                        # Slide title
                        title = slide.placeholders[current_slide_format_json["title_placeholder"]]
                        title.text = validated_slide_json_content["title"]
                        # Slide subtitle
                        subtitle = slide.placeholders[current_slide_format_json["subtitle_placeholder"]]
                        subtitle.text = validated_slide_json_content["subtitle"]
                        # Slide main text
                        main_text = slide.placeholders[current_slide_format_json["text_placeholder"]]
                        main_text.text = validated_slide_json_content.get("text").replace("*** ","\n").replace("- ","").rstrip().lstrip()

                    elif current_slide_format == "Slide with image and text":
                        # Slide title
                        title = slide.placeholders[current_slide_format_json["title_placeholder"]]
                        title.text = validated_slide_json_content["title"]
                        # Slide main text
                        main_text = slide.placeholders[current_slide_format_json["text_placeholder"]]
                        main_text.text = validated_slide_json_content.get("text").replace("*** ","\n").replace("- ","").rstrip().lstrip()
                        # Slide image
                        if st.session_state["generate_images"]:
                            image_placeholder = slide.placeholders[current_slide_format_json["image_placeholder"]]
                            summary_prompt = """Summarize the following content in comma separated abstract concepts, maximum 20 words. 
            The text will be used to generate a representative image with Stable Diffusion. Remove preamble when answering. Content:\n"""+(main_text.text if validated_slide_json_content["slideFormat"] == "Slide with image and text" else validated_slide_json_content["title"])
                            summary_prompt2, usage = generate_text(prompt=summary_prompt, model_id = st.session_state["chosen_LLM"])
                            st.session_state["n_input_tokens"] = st.session_state["n_input_tokens"] + usage['input_tokens']
                            st.session_state["n_output_tokens"] = st.session_state["n_output_tokens"] + usage['output_tokens']
                            
                            print("summary_prompt for image generation: "+str(summary_prompt2))
                            img_prompt = summary_prompt2[0].get("text")+ ", abstract"

                            generate_bedrock_image(img_prompt=img_prompt, current_slide_format_json=current_slide_format_json, image_placeholder=image_placeholder, cwd=cwd)
                            st.session_state["n_gen_images"] = st.session_state["n_gen_images"]+1

                    elif current_slide_format == "Slide with image only":
                        # Slide image
                        if st.session_state["generate_images"]:
                            image_placeholder = slide.placeholders[current_slide_format_json["image_placeholder"]]
                            summary_prompt = """Summarize the following content in comma separated abstract concepts, maximum 20 words. 
            The text will be used to generate a representative image with Stable Diffusion. Remove preamble when anwering. Content:\n"""+(main_text.text if validated_slide_json_content["slideFormat"] == "Slide with image and text" else validated_slide_json_content["title"])
                            summary_prompt2, usage = generate_text(prompt=summary_prompt, model_id = st.session_state["chosen_LLM"])
                            st.session_state["n_input_tokens"] = st.session_state["n_input_tokens"] + usage['input_tokens']
                            st.session_state["n_output_tokens"] = st.session_state["n_output_tokens"] + usage['output_tokens']
                            
                            print("summary_prompt for image generation: "+str(summary_prompt2))
                            img_prompt = summary_prompt2[0].get("text")+ ", abstract"

                            generate_bedrock_image(img_prompt=img_prompt, current_slide_format_json=current_slide_format_json, image_placeholder=image_placeholder, cwd=cwd)
                            st.session_state["n_gen_images"] = st.session_state["n_gen_images"]+1

                    elif current_slide_format == "Slide with 4 takeaways":
                        # Slide title
                        title = slide.placeholders[current_slide_format_json["title_placeholder"]]
                        title.text = validated_slide_json_content["title"]
                        # Slide 4 key takeaways
                        four_options = validated_slide_json_content.get("text").split("***")
                        four_options = filter(None, four_options)
                        for i, text_opt in enumerate(four_options):
                            if text_opt.rstrip().lstrip() == "": 
                                continue
                            if i>3: 
                                continue
                            text_opt_placeholder = slide.placeholders[current_slide_format_json["text"+str(i+1)+"_placeholder"]]
                            print(str(i+1)+" text_opt_placeholder.text before "+text_opt_placeholder.text)
                            text_opt_placeholder.text = text_opt.rstrip().lstrip()
                            print(str(i+1)+" text_opt_placeholder.text "+text_opt_placeholder.text)
                        
                if current_slide_format == "Title page" and st.session_state["create_agenda_checkbox"]:
                    # If required, create an Agenda after Title slide
                    i_slide_col = i_slide_col+1
                    outputs_col1[i_slide_col], outputs_col2[i_slide_col], outputs_col3[i_slide_col] = st.columns(3)
                    with outputs_col1[i_slide_col]:
                        # Pick slide layout from master
                        current_slide_format = "Agenda"
                        st.write("Slide "+str(i_slide_col+1)+" - format: "+current_slide_format)
                        current_slide_format_json = st.session_state["slides_format_json"][current_slide_format]
                        layout_slide_idx = current_slide_format_json["layout_slide"]
                        slide_layout = prs.slide_layouts[layout_slide_idx]
                        slide = prs.slides.add_slide(slide_layout)
                        # Slide title
                        title = slide.placeholders[current_slide_format_json["title_placeholder"]]
                        title.text = current_slide_format_json["agenda_title"]
                        agenda_items = []
                        agenda_slide_json_fix_attempts = 0
                        max_agenda_slide_json_fix_attempts = 3
                        while type(agenda_items) != "str" and agenda_slide_json_fix_attempts < max_agenda_slide_json_fix_attempts:
                            agenda_slide_json_fix_attempts = agenda_slide_json_fix_attempts+1
                            result_agenda_items, usage = invoke_llm_text(prompt= agenda_prompt(SLIDE_TITLES=[item["title"] for item in json_content_list]), model_id = st.session_state["chosen_LLM"])
                            st.session_state["n_input_tokens"] = st.session_state["n_input_tokens"] + usage['input_tokens']
                            st.session_state["n_output_tokens"] = st.session_state["n_output_tokens"] + usage['output_tokens']
                            try:
                                agenda_items = ast.literal_eval( (result_agenda_items[0])["text"] )["agenda_points"]
                            except ValueError:
                                print("FAILED TO GENERATE AGENDA")
                                agenda_items = []

                        # Slide main text
                        main_text = slide.placeholders[current_slide_format_json["text_placeholder"]]
                        try:
                            agenda_items_text = agenda_items.replace("*** ","\n").replace("- ","").rstrip().lstrip()
                        except ValueError:
                            agenda_items_text = str(agenda_items).replace("[","").replace("]","")
                        if len(agenda_items_text)==1:
                            agenda_items_text = agenda_items_text.replace(", ","\n")
                        st.write(agenda_items_text)
                        main_text.text = agenda_items_text
                        # Check shapes
                        print("layout_slide_idx",layout_slide_idx)
                        print("current_slide_format",current_slide_format)
                        for i,shape in enumerate(slide.placeholders):
                            print(str('\tPlaceholder Index: %d - Name: %s' % (shape.placeholder_format.idx, shape.name)))

                if current_slide_format == "Slide with 4 takeaways" and st.session_state["create_thankyou_checkbox"]:
                    # If required, create a Thank You after takeaway slide
                    i_slide_col = i_slide_col+1
                    outputs_col1[i_slide_col], outputs_col2[i_slide_col], outputs_col3[i_slide_col] = st.columns(3)
                    with outputs_col1[i_slide_col]:
                        # Pick slide layout from master
                        current_slide_format = "Thank you"
                        st.write("Slide "+str(i_slide_col+1)+" - format: "+current_slide_format)
                        current_slide_format_json = st.session_state["slides_format_json"][current_slide_format]
                        layout_slide_idx = current_slide_format_json["layout_slide"]
                        slide_layout = prs.slide_layouts[layout_slide_idx]
                        slide = prs.slides.add_slide(slide_layout)
                        
                        # Add Thank you! to the slide
                        thank_you = slide.placeholders[st.session_state["slides_format_json"]["Thank you"]["title_placeholder"]] 
                        thank_you.text = "Thank you!"

                        # Slide Author Name
                        if st.session_state["customize_contact_info"]:
                            try:
                                full_name = slide.placeholders[current_slide_format_json["full_name_placeholder"]] 
                                full_name.text = st.session_state["your_full_name"]
                            except ValueError:
                                print("FAILED TO RETRIEVE THE FULL NAME BOX!")
                        # Slide Contact
                        if st.session_state["customize_contact_info"]:
                            try:
                                contact_text = slide.placeholders[current_slide_format_json["contact_info_placeholder"]]
                                contact_text.text = st.session_state["your_contact_info"]
                            except ValueError:
                                print("FAILED TO RETRIEVE THE CONTACT INFO BOX!")

                        # Check shapes
                        print("layout_slide_idx",layout_slide_idx)
                        print("current_slide_format",current_slide_format)
                        for i,shape in enumerate(slide.placeholders):
                            print(str('\tPlaceholder Index: %d - Name: %s' % (shape.placeholder_format.idx, shape.name)))

                        left = top = util.Inches(0)
                        if st.session_state["selected_generate_bkg"]:
                            bkg_img_path = "/tmp/test_image_bkg.jpg"
                        else:
                            bkg_img_path = "/templates/default_bkg.jpg"
                        pic = slide.shapes.add_picture(cwd+bkg_img_path, left, top, width=prs.slide_width, height=prs.slide_height)
                        cursor_sp = slide.shapes[0]._element
                        cursor_sp.addprevious(pic._element)


                i_slide_col = i_slide_col+1
                st.write("")

            # Save presentation
            st.session_state["output_file"] = cwd+'/output/output_'+str(uuid.uuid4())+'.pptx'
            prs.save(st.session_state["output_file"])

            if st.session_state["generate_thumbnails"]:
                # GENERATE AND DISPLAY THUMBNAILS
                st.session_state["output_dir_images"] = st.session_state["output_file"].replace('.pptx',"/")
                _cmd = str('unoconv -o '+st.session_state["output_dir_images"]+'  -f html '+st.session_state["output_file"])
                subprocess.run(shlex.split(_cmd), shell=False) # nosemgrep: dangerous-subprocess-use-audit, input not controllable by an external resource / no user input
                images = glob.glob(st.session_state["output_dir_images"]+"*.jpg")
                for index in range(len(images)):
                    image = Image.open(images[index])
                    with outputs_col2[index]:
                        st.image(image, use_column_width=True)
            
            stop_time = datetime.now()
            delta = stop_time - start_time
            input_token_cost_cent = st.session_state["n_input_tokens"]*st.session_state["LLM_input_token_price"]*100
            output_token_cost_cent = st.session_state["n_output_tokens"]*st.session_state["LLM_output_token_price"]*100
            image_cost_cent = st.session_state["n_gen_images"]*st.session_state["gen_images_price_cents"]
            total_cost_cent = input_token_cost_cent+output_token_cost_cent+image_cost_cent
            st.warning(('''Generation completed in '''+str(round(float(delta.total_seconds())))+''' seconds.  
    '''+('''  - LLM input tokens: '''+str(st.session_state["n_input_tokens"]))+(''' (¢ '''+str(round(input_token_cost_cent,3))+''')  
    - LLM output tokens: '''+str(st.session_state["n_output_tokens"]))+(''' (¢ '''+str(round(output_token_cost_cent,3)))+''')  
    - Generated images: '''+str(st.session_state["n_gen_images"]))+(''' (¢ '''+str(round(image_cost_cent,3)))+''')  
      
    Total cost in cents of $: ¢ '''+str(round(total_cost_cent,2)), icon="⚠️")
            st.write("")
            with open(st.session_state["output_file"], "rb") as file:
                btn = st.download_button(
                        label="Download generated presentation",
                        data=file,
                        file_name="your_generated_presentation.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
    else:
        st.write("Please choose another topic or try again\n")